"""Generate COS760 Group 57 presentation using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DARK_BLUE = RGBColor(0, 51, 102)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 245)
BLACK = RGBColor(0, 0, 0)
FONT = "Calibri"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_font(run, size=18, bold=False, color=BLACK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_bar(slide, title_text):
    """Add dark blue title bar at top of slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_top = Pt(12)
    tf.margin_left = Pt(24)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    set_font(run, size=28, bold=True, color=WHITE)


def add_bullets(slide, left, top, width, height, items, size=18):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"• {item}"
        set_font(run, size=size)
    return txBox


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return txBox


# Create presentation
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank_layout = prs.slide_layouts[6]  # Blank layout

# ============ SLIDE 1 — Title ============
slide = prs.slides.add_slide(blank_layout)
# Full dark blue background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()

add_text(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
         "Cross-Lingual Detection of\nMachine-Generated Text in isiZulu",
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
         "Using Transfer Learning and Explainability Analysis",
         size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
         "Group 57 | COS760 | University of Pretoria",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)

# ============ SLIDE 2 — Problem & Motivation ============
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Problem & Motivation")
bullets = [
    "LLMs now generate fluent text in 100+ languages",
    "isiZulu: 12M+ speakers, most spoken home language in South Africa",
    "Risks: academic fraud, misinformation, fake government communications",
    "Gap: Zero existing MGT detection research for isiZulu or Bantu languages",
    "Our contribution: First transformer-based MGT detector for isiZulu + explainability analysis",
]
add_bullets(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5), bullets, size=20)

# ============ SLIDE 3 — Research Questions ============
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Research Questions")
rq1 = "RQ1: How accurately can fine-tuned AfroXLMR detect machine-generated isiZulu text?"
rq2 = "RQ2: Does an English-trained classifier transfer to isiZulu, and what does LIME reveal about cross-lingual detection features?"
add_text(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5), rq1,
         size=22, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(1.5), rq2,
         size=22, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# ============ SLIDE 4 — Data Pipeline ============
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Data Pipeline")

# Left column
add_text(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
         "Human Text", size=20, bold=True, color=DARK_BLUE)
add_bullets(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5), [
    "isiZulu: Vukuzenzele + Wikipedia",
    "English: News articles + Wikipedia",
    "Chunked into 4-sentence segments (\u226530 words)",
], size=18)

# Right column
add_text(slide, Inches(6.8), Inches(1.5), Inches(5), Inches(0.5),
         "Machine Text", size=20, bold=True, color=DARK_BLUE)
add_bullets(slide, Inches(6.8), Inches(2.1), Inches(5.5), Inches(2.5), [
    "Generator: Gemma 4 (26B parameters)",
    "185 topic cues \u00d7 4 prompt templates",
    "Both isiZulu and English",
], size=18)

# Bottom section
add_text(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.5),
         "Splits: 70% train / 15% val / 15% test", size=16, bold=True)
add_text(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.8),
         "Configs: isiZulu-only (946/203/203) | English-only (977/209/210) | Combined | Cross-lingual (Eng\u2192Zul)",
         size=16)

# ============ SLIDE 5 — Models & Method ============
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Models & Method")

# Left column
add_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
         "Primary Model", size=20, bold=True, color=DARK_BLUE)
add_bullets(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.2), [
    "AfroXLMR-base (Davlan/afro-xlmr-base)",
    "278M params, 12 layers",
    "Additional pretraining on 17 African languages incl. isiZulu",
], size=18)

# Right column
add_text(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.5),
         "Baseline", size=20, bold=True, color=DARK_BLUE)
add_bullets(slide, Inches(6.8), Inches(2.1), Inches(5.5), Inches(2.2), [
    "XLM-RoBERTa-base",
    "278M params, same architecture",
    "General multilingual (100 languages, no African specialisation)",
], size=18)

# Bottom
add_text(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.5),
         "Training: lr=2\u00d710\u207b\u2075, batch=16, max_len=256, early stopping (patience=2), seed=42",
         size=16, bold=True)
add_text(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.5),
         "Explainability: LIME \u2014 top 12 tokens per sample, 10 samples per model", size=16)

# ============ SLIDE 6 — Results ============
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Results")

# Table data
headers = ["Model", "Config", "F1", "Accuracy", "ROC-AUC"]
rows_data = [
    ["AfroXLMR", "isiZulu", "0.975", "0.975", "0.995"],
    ["AfroXLMR", "English", "0.995", "0.995", "1.000"],
    ["AfroXLMR", "Combined", "0.985", "0.985", "0.995"],
    ["AfroXLMR", "Cross-lingual", "0.737", "0.754", "0.985"],
    ["XLM-R", "isiZulu", "0.941", "0.941", "0.993"],
    ["XLM-R", "English", "0.995", "0.995", "1.000"],
    ["XLM-R", "Combined", "0.983", "0.983", "0.995"],
    ["XLM-R", "Cross-lingual", "0.417", "0.542", "0.973"],
    ["Majority", "All", "0.334", "0.502", "\u2014"],
]

num_rows = len(rows_data) + 1
num_cols = len(headers)
tbl_left = Inches(0.5)
tbl_top = Inches(1.4)
tbl_width = Inches(12.3)
tbl_height = Inches(3.8)

table_shape = slide.shapes.add_table(num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_height)
table = table_shape.table

# Set column widths
col_widths = [Inches(2.2), Inches(2.8), Inches(2.2), Inches(2.2), Inches(2.2)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

# Header row
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = h
    set_font(run, size=14, bold=True, color=WHITE)

# Data rows
for r_idx, row in enumerate(rows_data):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx + 1, c_idx)
        # Alternating row shading
        if r_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = val
        set_font(run, size=13)

# Key findings below table
add_bullets(slide, Inches(0.5), Inches(5.4), Inches(12), Inches(2), [
    "AfroXLMR beats XLM-R on isiZulu by 3.4 F1 points",
    "Cross-lingual gap: 32 F1 points (AfroXLMR 0.737 vs XLM-R 0.417)",
    "Monolingual detection near-perfect for both models",
], size=16)

# ============ SLIDE 7 — LIME Explainability ============
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "LIME Explainability Analysis")

# Left column
add_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
         "isiZulu Model (strong signal)", size=18, bold=True, color=DARK_BLUE)
add_bullets(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5), [
    "Top tokens: kubalulekile, okuhloselwe, Ngokusebenzisana",
    "Avg weight: 0.13\u20130.23",
    "Detects: formal vocabulary, discourse markers overused by LLMs",
], size=16)

# Right column
add_text(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.5),
         "Cross-lingual Model (weak signal)", size=18, bold=True, color=DARK_BLUE)
add_bullets(slide, Inches(6.8), Inches(2.1), Inches(5.5), Inches(2.5), [
    "Top tokens: Kubalulekile, ezitholakala, izindleko",
    "Avg weight: 0.02\u20130.03 (10\u00d7 weaker)",
    "Near-random feature reliance",
], size=16)

# Bottom
add_bullets(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2), [
    "Jaccard similarity of top features: 0.043 (almost no overlap)",
    "Conclusion: English detection patterns do NOT transfer to isiZulu morphology",
], size=16)

# ============ SLIDE 8 — Conclusion & Future Work ============
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Conclusion & Future Work")

add_text(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
         "Conclusion", size=20, bold=True, color=DARK_BLUE)
add_bullets(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.5), [
    "AfroXLMR achieves F1 = 0.975 on isiZulu MGT detection \u2713",
    "Cross-lingual transfer is feasible but limited (F1 = 0.737) \u2014 language-specific training essential",
    "LIME confirms fundamentally different detection strategies per language",
    "African language pretraining provides critical advantage (+3.4 F1 mono, +32 F1 cross)",
], size=17)

add_text(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.5),
         "Future Work", size=20, bold=True, color=DARK_BLUE)
add_bullets(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2), [
    "Multiple generators (GPT-4, Claude, Llama)",
    "More African languages (Xhosa, Sotho, Swahili)",
    "Threshold calibration for cross-lingual deployment",
    "Few-shot adaptation as middle ground",
], size=17)

# Save
output_path = r"C:\Users\reece\cos\760\groupwork\final_pres\presentation\COS760_Group57_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
